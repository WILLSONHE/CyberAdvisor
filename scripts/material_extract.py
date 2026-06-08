#!/usr/bin/env python3
"""
从 Raw/未分析其他材料/ 提取文本，供 ing 消化。

用法:
  python material_extract.py list
  python material_extract.py extract FILE
  python material_extract.py extract-all [--out-dir DIR]
"""
from __future__ import annotations

import argparse
import os
import sys

from raw_paths import RAW_PENDING_MATERIALS, WIKI_MATERIALS, list_pending_material_files, ensure_raw_dirs


def _read_text_file(path: str) -> str:
    for enc in ("utf-8", "utf-8-sig", "gbk", "gb18030"):
        try:
            with open(path, encoding=enc) as f:
                return f.read()
        except UnicodeDecodeError:
            continue
    with open(path, encoding="utf-8", errors="replace") as f:
        return f.read()


def _extract_pdf(path: str) -> str:
    try:
        from pypdf import PdfReader
    except ImportError as e:
        raise RuntimeError("需要安装 pypdf：pip install pypdf") from e
    reader = PdfReader(path)
    parts: list[str] = []
    for i, page in enumerate(reader.pages, start=1):
        text = (page.extract_text() or "").strip()
        if text:
            parts.append(f"--- 第 {i} 页 ---\n{text}")
    return "\n\n".join(parts)


def _extract_pptx(path: str) -> str:
    try:
        from pptx import Presentation
    except ImportError as e:
        raise RuntimeError("需要安装 python-pptx：pip install python-pptx") from e
    prs = Presentation(path)
    parts: list[str] = []
    for i, slide in enumerate(prs.slides, start=1):
        lines: list[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                lines.append(shape.text.strip())
        if lines:
            parts.append(f"--- 幻灯片 {i} ---\n" + "\n".join(lines))
    return "\n\n".join(parts)


def _extract_docx(path: str) -> str:
    try:
        from docx import Document
    except ImportError as e:
        raise RuntimeError("需要安装 python-docx：pip install python-docx") from e
    doc = Document(path)
    return "\n\n".join(p.text.strip() for p in doc.paragraphs if p.text.strip())


def extract_text(path: str) -> tuple[str, str]:
    """返回 (格式, 正文)。"""
    ext = os.path.splitext(path)[1].lower()
    if ext in (".md", ".txt"):
        return ext.lstrip("."), _read_text_file(path)
    if ext == ".pdf":
        return "pdf", _extract_pdf(path)
    if ext == ".pptx":
        return "pptx", _extract_pptx(path)
    if ext == ".docx":
        return "docx", _extract_docx(path)
    raise ValueError(f"不支持的格式: {ext}")


def cmd_list() -> int:
    ensure_raw_dirs()
    files = list_pending_material_files()
    if not files:
        print(f"（空）{RAW_PENDING_MATERIALS}")
        return 0
    for p in files:
        ext = os.path.splitext(p)[1].lower()
        print(f"{ext}\t{p}")
    print(f"\n共 {len(files)} 个文件")
    return 0


def cmd_extract(path: str, *, out: str = "") -> int:
    if not os.path.isfile(path):
        raise SystemExit(f"文件不存在: {path}")
    fmt, text = extract_text(path)
    if not text.strip():
        print(f"[WARN] {path} 提取结果为空（{fmt}）", file=sys.stderr)
    header = f"# 来源: {os.path.basename(path)}\n\n> 格式: {fmt}\n\n"
    body = header + text.strip() + "\n"
    if out:
        os.makedirs(os.path.dirname(os.path.abspath(out)), exist_ok=True)
        with open(out, "w", encoding="utf-8") as f:
            f.write(body)
        print(out)
    else:
        print(body)
    return 0


def cmd_extract_all(*, out_dir: str) -> int:
    ensure_raw_dirs()
    os.makedirs(out_dir, exist_ok=True)
    files = list_pending_material_files()
    if not files:
        print("无待提取文件")
        return 0
    for src in files:
        stem = os.path.splitext(os.path.basename(src))[0]
        dest = os.path.join(out_dir, f"{stem}.md")
        cmd_extract(src, out=dest)
    print(f"已提取 {len(files)} 个 → {out_dir}")
    return 0


def main() -> None:
    ap = argparse.ArgumentParser(description="其他材料文本提取（ing 前置）")
    sub = ap.add_subparsers(dest="cmd", required=True)
    sub.add_parser("list", help="列出未分析其他材料")
    ex = sub.add_parser("extract", help="提取单个文件")
    ex.add_argument("file", help="文件路径")
    ex.add_argument("--out", default="", help="写入 md 路径（默认 stdout）")
    ea = sub.add_parser("extract-all", help="批量提取到目录")
    ea.add_argument(
        "--out-dir",
        default=os.path.join(WIKI_MATERIALS, "_extracts"),
        help="输出目录（默认 Wiki/其他材料/_extracts）",
    )
    args = ap.parse_args()
    if args.cmd == "list":
        raise SystemExit(cmd_list())
    if args.cmd == "extract":
        raise SystemExit(cmd_extract(args.file, out=args.out))
    if args.cmd == "extract-all":
        raise SystemExit(cmd_extract_all(out_dir=args.out_dir))


if __name__ == "__main__":
    main()
